from dataclasses import dataclass
from collections import namedtuple
from typing import Dict
from pptx import Presentation

from presentation.mySlide import MySlide


@dataclass
class MyPresentation:
    """
    Class representing a PowerPoint presentation.
    """

    filepath: str
    explanations: namedtuple = None
    slides: Dict[int, MySlide] = None

    def __post_init__(self):
        """
        Initializes the slides dictionary.
        """
        self.slides = {}

    def parse(self):
        """
        Parses the PowerPoint presentation to extract slide data.
        """
        prs = Presentation(self.filepath)
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_obj = prs.slides[slide_number - 1]  # Access slide object using slide number - 1

            # Extract text from text boxes on the slide
            slide_text_boxes = [
                run.text.strip()
                for shape in slide_obj.shapes
                if shape.has_text_frame
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.text.strip()
            ]

            # Create a MySlide object and store it in the slides dictionary
            self.slides[slide_number] = MySlide(slide_number, slide_text_boxes)
